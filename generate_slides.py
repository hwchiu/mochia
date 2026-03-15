from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Theme ─────────────────────────────────────────────────────────────────────
BG_COLOR     = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy background
ACCENT_BLUE  = RGBColor(0x00, 0xBF, 0xFF)   # electric blue
ACCENT_GREEN = RGBColor(0x39, 0xFF, 0x14)   # neon green
ACCENT_AMBER = RGBColor(0xFF, 0xBF, 0x00)   # amber
TEXT_PRIMARY = RGBColor(0xE8, 0xED, 0xF2)   # off-white
TEXT_DIM     = RGBColor(0x8F, 0xA3, 0xBF)   # dimmed subtext
CARD_BG      = RGBColor(0x1E, 0x2D, 0x3D)   # card background
CARD_BORDER  = RGBColor(0x00, 0x5F, 0x8A)   # card border

FONT_TITLE = "JetBrains Mono"
FONT_BODY  = "Inter"
SLIDE_W    = Inches(13.33)
SLIDE_H    = Inches(7.5)

PART_COLORS = {
    1: ACCENT_BLUE,
    2: RGBColor(0x7B, 0x2F, 0xFF),
    3: RGBColor(0x00, 0xE5, 0xFF),
    4: RGBColor(0xFF, 0x6B, 0x35),
    5: ACCENT_GREEN,
    6: ACCENT_AMBER,
}


def new_slide(prs, layout_idx=6):
    """Add blank slide and paint background."""
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    return slide


def add_rect(slide, l, t, w, h, fill_color, border_color=None, border_width=Pt(1)):
    """Add a filled rectangle shape. MSO_SHAPE_TYPE 1 = rectangle."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, font_name=FONT_BODY, font_size=Pt(14),
             color=TEXT_PRIMARY, bold=False, align=PP_ALIGN.LEFT, italic=False):
    """Add a text box with a single paragraph."""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return txBox


def add_multiline_text(slide, lines, l, t, w, h, font_name=FONT_BODY,
                       font_size=Pt(13), color=TEXT_PRIMARY):
    """Add a textbox with multiple paragraphs.
    lines: list of (text, opts_dict) where opts can have: align, font, size, color, bold, italic
    """
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, opts) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = opts.get("align", PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = text
        run.font.name = opts.get("font", font_name)
        run.font.size = opts.get("size", font_size)
        run.font.color.rgb = opts.get("color", color)
        run.font.bold = opts.get("bold", False)
        run.font.italic = opts.get("italic", False)
    return txBox


def add_part_label(slide, part_num, part_title):
    """Top-left Part badge with colored background."""
    color = PART_COLORS.get(part_num, ACCENT_BLUE)
    add_rect(slide, Inches(0.3), Inches(0.15), Inches(1.5), Inches(0.32), color)
    add_text(slide, f"PART {part_num}", Inches(0.3), Inches(0.15), Inches(1.5), Inches(0.32),
             font_name=FONT_TITLE, font_size=Pt(9),
             color=RGBColor(0x0D, 0x1B, 0x2A), bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, part_title, Inches(1.9), Inches(0.18), Inches(8), Inches(0.28),
             font_name=FONT_BODY, font_size=Pt(9), color=TEXT_DIM)


def add_page_number(slide, num):
    """Bottom-right page number."""
    add_text(slide, f"{num:02d} / 50", Inches(12.0), Inches(7.1), Inches(1.1), Inches(0.28),
             font_size=Pt(9), color=TEXT_DIM, align=PP_ALIGN.RIGHT)


def add_slide_title(slide, title, subtitle=None):
    """Standard slide title with blue underline rule."""
    add_text(slide, title, Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.7),
             font_name=FONT_TITLE, font_size=Pt(26), color=ACCENT_BLUE, bold=True)
    add_rect(slide, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.025), ACCENT_BLUE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.5), Inches(1.3), Inches(10), Inches(0.35),
                 font_size=Pt(12), color=TEXT_DIM)


def add_card(slide, l, t, w, h, title=None, body_lines=None, accent=ACCENT_BLUE):
    """Dark card with colored border and optional title bar + body lines."""
    add_rect(slide, l, t, w, h, CARD_BG, accent, Pt(1.5))
    y_offset = t
    if title:
        add_rect(slide, l, t, w, Inches(0.35), accent)
        add_text(slide, title, l + Inches(0.1), t + Inches(0.03),
                 w - Inches(0.2), Inches(0.3),
                 font_name=FONT_TITLE, font_size=Pt(10),
                 color=RGBColor(0x0D, 0x1B, 0x2A), bold=True)
        y_offset = t + Inches(0.42)
    else:
        y_offset = t + Inches(0.1)

    if body_lines:
        y = y_offset
        for line in body_lines:
            add_text(slide, line, l + Inches(0.12), y,
                     w - Inches(0.24), Inches(0.28),
                     font_size=Pt(11), color=TEXT_PRIMARY)
            y += Inches(0.3)


def add_callout(slide, text, l, t, w, h, style="tip"):
    """Callout box. style: 'tip' (green), 'warning' (amber), 'danger' (red)."""
    configs = {
        "tip":     (RGBColor(0x00, 0x1A, 0x0A), ACCENT_GREEN),
        "warning": (RGBColor(0x1A, 0x10, 0x00), ACCENT_AMBER),
        "danger":  (RGBColor(0x1A, 0x00, 0x00), RGBColor(0xFF, 0x4A, 0x4A)),
    }
    bg_c, border_c = configs.get(style, configs["tip"])
    add_rect(slide, l, t, w, h, bg_c, border_c, Pt(1.5))
    prefix = {"tip": "💡  ", "warning": "⚠   ", "danger": "❌  "}.get(style, "")
    add_text(slide, prefix + text, l + Inches(0.12), t + Inches(0.08),
             w - Inches(0.24), h - Inches(0.16),
             font_size=Pt(12), color=TEXT_PRIMARY)


def slide_01_cover(prs):
    slide = new_slide(prs)

    # Full-width top accent line
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

    # Decorative vertical accent bar on left
    add_rect(slide, Inches(0), Inches(0.06), Inches(0.06), SLIDE_H - Inches(0.06), ACCENT_BLUE)

    # Background grid pattern (subtle horizontal lines)
    for i in range(8):
        y = Inches(1.0 + i * 0.8)
        add_rect(slide, Inches(0.3), y, SLIDE_W - Inches(0.3), Inches(0.005),
                 RGBColor(0x1E, 0x2D, 0x3D))

    # Main title: CLOUD NATIVE
    add_text(slide, "CLOUD NATIVE", Inches(1.2), Inches(1.5), Inches(11), Inches(1.4),
             font_name=FONT_TITLE, font_size=Pt(64), color=ACCENT_BLUE, bold=True,
             align=PP_ALIGN.LEFT)

    # Subtitle
    add_text(slide, "系統部署實務", Inches(1.2), Inches(2.9), Inches(10), Inches(0.85),
             font_name=FONT_TITLE, font_size=Pt(36), color=TEXT_PRIMARY, bold=False,
             align=PP_ALIGN.LEFT)

    # Underline rule
    add_rect(slide, Inches(1.2), Inches(3.8), Inches(10), Inches(0.04), ACCENT_BLUE)

    # Tagline
    add_text(slide, "從單體部署到 Cloud Native 的完整演進之路",
             Inches(1.2), Inches(3.9), Inches(10), Inches(0.45),
             font_size=Pt(16), color=TEXT_DIM, align=PP_ALIGN.LEFT)

    # Bottom info bar background
    add_rect(slide, Inches(0), Inches(6.5), SLIDE_W, Inches(1.0), CARD_BG)

    # Bottom left: course info
    add_text(slide, "碩士課程  ·  系統架構設計",
             Inches(0.5), Inches(6.62), Inches(6), Inches(0.35),
             font_size=Pt(13), color=TEXT_DIM)

    # Bottom right: stats
    add_text(slide, "預計課程時間：2.5 小時  ·  投影片：50 頁",
             Inches(7.0), Inches(6.62), Inches(6.0), Inches(0.35),
             font_size=Pt(13), color=TEXT_DIM, align=PP_ALIGN.RIGHT)

    # Small decorative dot cluster (top right)
    for row in range(4):
        for col in range(4):
            x = Inches(11.5) + col * Inches(0.22)
            y = Inches(1.0) + row * Inches(0.22)
            sz = Inches(0.08)
            add_rect(slide, x, y, sz, sz, CARD_BORDER)

    return slide


def slide_02_agenda(prs):
    slide = new_slide(prs)

    # Title
    add_slide_title(slide, "課程大綱", "Agenda — 完整演進之路")
    add_page_number(slide, 2)

    # 6 Part cards in 2 rows × 3 columns
    parts = [
        (1, "PART 1", "傳統部署演進",    "單機 → 三層架構，了解部署的起點"),
        (2, "PART 2", "Scale Out 挑戰",  "LB / Session / DB 擴展的真實難題"),
        (3, "PART 3", "Container 革命",  "Docker / Compose / Registry"),
        (4, "PART 4", "12-Factor App",   "Cloud-Ready 應用的設計原則"),
        (5, "PART 5", "DevOps 整合",     "大規模生產環境的團隊協作"),
        (6, "PART 6", "SDLC 閉環",       "從 Code 到維運的完整旅程"),
    ]

    card_w = Inches(4.0)
    card_h = Inches(1.9)
    x_starts = [Inches(0.4), Inches(4.6), Inches(8.8)]
    y_starts  = [Inches(1.55), Inches(3.65)]

    for i, (part_num, label, title, desc) in enumerate(parts):
        col = i % 3
        row = i // 3
        x = x_starts[col]
        y = y_starts[row]
        color = PART_COLORS[part_num]

        # Card background
        add_rect(slide, x, y, card_w, card_h, CARD_BG, color, Pt(1.5))

        # Top colored header strip
        add_rect(slide, x, y, card_w, Inches(0.4), color)

        # Part label in header
        add_text(slide, label, x + Inches(0.1), y + Inches(0.05),
                 card_w - Inches(0.2), Inches(0.32),
                 font_name=FONT_TITLE, font_size=Pt(11),
                 color=RGBColor(0x0D, 0x1B, 0x2A), bold=True)

        # Title
        add_text(slide, title, x + Inches(0.12), y + Inches(0.46),
                 card_w - Inches(0.24), Inches(0.4),
                 font_name=FONT_TITLE, font_size=Pt(14),
                 color=TEXT_PRIMARY, bold=True)

        # Description
        add_text(slide, desc, x + Inches(0.12), y + Inches(0.92),
                 card_w - Inches(0.24), Inches(0.85),
                 font_size=Pt(11), color=TEXT_DIM)

    # Bottom note
    add_text(slide, "課程目標：理解真實世界如何從零設計可擴展系統，體會各種設計決策背後的取捨",
             Inches(0.4), Inches(5.7), Inches(12.5), Inches(0.4),
             font_size=Pt(11), color=TEXT_DIM, align=PP_ALIGN.CENTER)

    return slide


def slide_03(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "起點：最簡單的部署架構")
    add_part_label(slide, 1, "傳統部署的演進")
    add_page_number(slide, 3)

    # Left side: Architecture diagram
    add_card(slide, Inches(0.5), Inches(1.8), Inches(5.5), Inches(0.7),
             title="使用者 Browser", accent=ACCENT_BLUE)
    add_text(slide, "→", Inches(0.5), Inches(2.55), Inches(5.5), Inches(0.4),
             font_size=Pt(22), color=TEXT_DIM, align=PP_ALIGN.CENTER, bold=True)
    add_card(slide, Inches(0.5), Inches(3.0), Inches(5.5), Inches(0.9),
             title="Linux Server :8080 (my-app)", accent=ACCENT_BLUE)

    # Right side: Two cards stacked
    add_card(slide, Inches(6.5), Inches(1.8), Inches(6.3), Inches(1.5),
             title="✅ 優點", accent=ACCENT_GREEN,
             body_lines=["設定簡單、維運容易", "適合初期開發與 PoC 驗證", "開發環境 ≈ 生產環境"])
    add_card(slide, Inches(6.5), Inches(3.5), Inches(6.3), Inches(1.5),
             title="⚠  問題", accent=ACCENT_AMBER,
             body_lines=["單點故障 (SPOF)", "無法應付高流量", "資料存記憶體 → 重啟即遺失"])

    # Callout tip
    add_callout(slide, "課堂 PoC 常用這種架構，但真實產品幾乎不會停在這裡",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_04(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "請求的旅程：一個 HTTP Request 經歷什麼？")
    add_part_label(slide, 1, "傳統部署的演進")
    add_page_number(slide, 4)

    steps = [
        ("① DNS 解析", "查詢域名 → 取得 Server IP"),
        ("② TCP 握手", "三次握手 → 建立可靠連線"),
        ("③ HTTP Request", "GET /api/data → 發送請求"),
        ("④ App 處理", "業務邏輯 → 查 DB / 計算"),
        ("⑤ DB Query", "執行 SQL → 取得資料"),
        ("⑥ HTTP Response", "JSON 回傳 → Client 收到"),
    ]

    x_positions = [Inches(0.4), Inches(6.9)]
    y_start = Inches(1.6)
    y_spacing = Inches(1.7)

    for i, (title, desc) in enumerate(steps):
        col = i % 2
        row = i // 2
        x = x_positions[col]
        y = y_start + row * y_spacing
        add_card(slide, x, y, Inches(5.8), Inches(1.4),
                 title=title, body_lines=[desc], accent=ACCENT_BLUE)

    add_callout(slide, "理解這條路徑，才知道瓶頸藏在哪一層",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_05(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "第一步擴展：分離資料庫 Server")
    add_part_label(slide, 1, "傳統部署的演進")
    add_page_number(slide, 5)

    # Left side: Architecture diagram
    add_card(slide, Inches(0.5), Inches(2.0), Inches(5.5), Inches(0.9),
             title="App Server  (my-app :8080)", accent=ACCENT_BLUE)
    add_text(slide, "→ SQL →", Inches(0.5), Inches(2.95), Inches(5.5), Inches(0.4),
             font_size=Pt(16), color=TEXT_DIM, align=PP_ALIGN.CENTER, bold=True)
    add_card(slide, Inches(0.5), Inches(3.3), Inches(5.5), Inches(0.9),
             title="DB Server  (PostgreSQL :5432)", accent=PART_COLORS[4])

    # Right side
    add_card(slide, Inches(6.5), Inches(2.0), Inches(6.3), Inches(1.7),
             title="為什麼要把 DB 獨立出來？",
             body_lines=[
                 "① App Server 可以無狀態重啟，不影響資料",
                 "② DB 資源需求不同，可以獨立調整規格",
                 "③ 資料可以被多個 App 共用",
             ],
             accent=ACCENT_BLUE)

    add_callout(slide, "新問題：網路連線延遲增加、DB 成為效能瓶頸、兩台機器各別維護",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="warning")
    return slide


def slide_06(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "三層架構：Frontend + Backend + Database")
    add_part_label(slide, 1, "傳統部署的演進")
    add_page_number(slide, 6)

    tiers = [
        (Inches(0.4),  PART_COLORS[3], "Presentation Tier",
         ["Frontend Server", "React / Vue", "Nginx : 80", "", "職責：UI 呈現、", "靜態資源服務"]),
        (Inches(4.7),  ACCENT_BLUE,    "Application Tier",
         ["Backend Server", "FastAPI / Node.js", ": 8080", "", "職責：業務邏輯、", "API 處理"]),
        (Inches(9.0),  PART_COLORS[4], "Data Tier",
         ["Database Server", "PostgreSQL", ": 5432", "", "職責：資料持久化、", "查詢處理"]),
    ]

    for x, accent, title, body in tiers:
        add_card(slide, x, Inches(1.6), Inches(3.8), Inches(3.8),
                 title=title, body_lines=body, accent=accent)

    # Arrows between cards
    add_text(slide, "→", Inches(4.2), Inches(3.0), Inches(0.5), Inches(0.5),
             font_size=Pt(22), color=TEXT_DIM, align=PP_ALIGN.CENTER, bold=True)
    add_text(slide, "→", Inches(8.5), Inches(3.0), Inches(0.5), Inches(0.5),
             font_size=Pt(22), color=TEXT_DIM, align=PP_ALIGN.CENTER, bold=True)

    add_callout(slide, "3 台機器 = 3 倍維運工作量，且任何一台掛掉都影響整個服務",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="warning")
    return slide


def slide_07(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "三層架構的真實挑戰")
    add_part_label(slide, 1, "傳統部署的演進")
    add_page_number(slide, 7)

    danger_color = RGBColor(0xFF, 0x4A, 0x4A)
    cards = [
        (Inches(0.4), Inches(1.6), "① 部署順序地雷",
         ["前後端必須依序更新", "Frontend v2 + Backend v1 = API 不相容", "任一步驟失敗，整體需回滾"]),
        (Inches(6.9), Inches(1.6), "② 版本相依性地獄",
         ["Frontend 假設 API 回傳格式", "Backend 改動 response 結構", "雙方沒有明確合約 → 上線即爆炸"]),
        (Inches(0.4), Inches(3.8), "③ 開發環境差異",
         ["本地 macOS → 測試 Ubuntu 20 → 生產 Ubuntu 22", "套件版本不同 → 行為不一致", "「我這裡跑得好好的...」"]),
        (Inches(6.9), Inches(3.8), "④ 單點故障",
         ["任何一層掛掉 → 整個服務中斷", "Frontend 掛 → 用戶看不到頁面", "DB 掛 → 全站 500 Error"]),
    ]

    for x, y, title, body in cards:
        add_card(slide, x, y, Inches(5.8), Inches(2.0),
                 title=title, body_lines=body, accent=danger_color)

    add_callout(slide, "這些問題在小團隊可以忍受，但隨著規模成長會爆炸",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="warning")
    return slide


def slide_08(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "如何找出系統瓶頸？")
    add_part_label(slide, 1, "傳統部署的演進")
    add_page_number(slide, 8)

    bottlenecks = [
        ("CPU 瓶頸",      "症狀: top/htop CPU 100%  |  工具: perf, flame graph"),
        ("Memory 瓶頸",   "症狀: OOM Kill, Swap 滿  |  工具: free, valgrind"),
        ("Disk I/O 瓶頸", "症狀: iowait 高, 寫入慢  |  工具: iostat, iotop"),
        ("Network 瓶頸",  "症狀: 封包遺失, 延遲高  |  工具: netstat, tcpdump"),
        ("DB Query 瓶頸", "症狀: 慢查詢, Lock 等待  |  工具: EXPLAIN ANALYZE"),
    ]

    y = Inches(1.6)
    for title, desc in bottlenecks:
        add_card(slide, Inches(0.4), y, Inches(7.5), Inches(0.85),
                 title=title, body_lines=[desc], accent=ACCENT_BLUE)
        y += Inches(0.9)

    # Right card
    add_card(slide, Inches(8.2), Inches(1.6), Inches(4.8), Inches(4.5),
             title="解讀監控訊號",
             body_lines=[
                 "P50 / P95 / P99 Latency",
                 "Request per Second (RPS)",
                 "Error Rate (%)",
                 "CPU / Memory 使用率",
                 "DB Connection Pool 使用量",
             ],
             accent=ACCENT_BLUE)

    add_callout(slide, "先量測，再優化。不要猜，要看數據。沒有監控就沒有優化的依據",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_09(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "何時需要開始思考 Scale？")
    add_part_label(slide, 1, "傳統部署的演進")
    add_page_number(slide, 9)

    indicators = [
        (Inches(0.4), RGBColor(0xFF, 0x4A, 0x4A), "CPU 持續 > 70%",
         ["計算資源不足", "壓力測試時先觸頂", "Scale Up 或 Scale Out", "加 Cache 降低計算量"]),
        (Inches(4.6), ACCENT_AMBER, "Response P99 > 1s",
         ["用戶體驗明顯下降", "長尾延遲影響轉換率", "查慢查詢 / 優化邏輯", "若已優化 → 需 Scale"]),
        (Inches(8.8), ACCENT_BLUE, "Error Rate > 0.1%",
         ["服務不穩定訊號", "可能是 OOM / DB 連線耗盡", "先找根因，再決定是否 Scale"]),
    ]

    for x, accent, title, body in indicators:
        add_card(slide, x, Inches(1.6), Inches(3.9), Inches(2.8),
                 title=title, body_lines=body, accent=accent)

    add_card(slide, Inches(0.4), Inches(4.65), Inches(12.5), Inches(1.3),
             title="常見錯誤：過早 Scale",
             body_lines=[
                 "需求還沒穩定就急著上 K8s",
                 "單機能搞定的問題別引入分散式複雜度",
                 "先用單機撐到真正需要的時候再 Scale",
             ],
             accent=ACCENT_AMBER)

    add_callout(slide, "過早優化是萬惡之源 — Donald Knuth。先讓產品活下來，再考慮 Scale",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_10(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "Scale Up vs Scale Out：兩種擴展思路")
    add_part_label(slide, 1, "傳統部署的演進")
    add_page_number(slide, 10)

    # Left: Scale Up
    add_card(slide, Inches(0.4), Inches(1.9), Inches(5.8), Inches(0.5),
             title="Scale Up  垂直擴展", accent=ACCENT_AMBER)
    add_card(slide, Inches(0.4), Inches(2.6), Inches(5.8), Inches(3.5),
             body_lines=[
                 "舊機器: 4 Core / 16GB",
                 "    ↓",
                 "新機器: 64 Core / 512GB",
                 "",
                 "✅ 架構不用改，快速",
                 "✅ 無需修改程式碼",
                 "⚠  成本指數級成長",
                 "⚠  硬體有上限",
                 "⚠  升級需要停機",
             ],
             accent=ACCENT_AMBER)

    # Right: Scale Out
    add_card(slide, Inches(7.0), Inches(1.9), Inches(5.8), Inches(0.5),
             title="Scale Out  水平擴展", accent=ACCENT_GREEN)
    add_card(slide, Inches(7.0), Inches(2.6), Inches(5.8), Inches(3.5),
             body_lines=[
                 "Server 1 (4C/16GB)",
                 "Server 2 (4C/16GB)",
                 "Server 3 (4C/16GB)  + N 台",
                 "",
                 "✅ 成本線性成長，彈性",
                 "✅ 零停機，滾動更新",
                 "✅ 理論上無上限",
                 "⚠  架構需要重新設計",
             ],
             accent=ACCENT_GREEN)

    # Center VS label
    add_text(slide, "VS", Inches(6.0), Inches(3.5), Inches(0.8), Inches(0.8),
             font_size=Pt(28), color=TEXT_DIM, align=PP_ALIGN.CENTER, bold=True)

    add_callout(slide, "Scale Out 是 Cloud Native 的核心，但需要應用程式支援 Stateless 設計",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_11(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "Stateless 設計：Scale Out 的先決條件")
    add_part_label(slide, 1, "傳統部署的演進")
    add_page_number(slide, 11)

    # Left: Stateful (bad)
    add_card(slide, Inches(0.4), Inches(1.6), Inches(5.8), Inches(3.5),
             title="❌ Stateful（有問題）",
             body_lines=[
                 "Session 存在 App Server 記憶體",
                 "Request #1 → Server A (登入成功)",
                 "Request #2 → Server B (找不到 Session!)",
                 "→ 用戶被強制登出！",
                 "",
                 "問題：Server 之間 Session 無法共用",
             ],
             accent=RGBColor(0xFF, 0x4A, 0x4A))

    # Right: Stateless (good)
    add_card(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(3.5),
             title="✅ Stateless（正確）",
             body_lines=[
                 "Session 外部化存入 Redis",
                 "Request #1 → Server A → Redis 取 Session",
                 "Request #2 → Server B → Redis 取 Session",
                 "→ 任一台 Server 都能處理！",
                 "",
                 "原則：App Server 不存任何特定狀態",
             ],
             accent=ACCENT_GREEN)

    # Center arrow
    add_text(slide, "→", Inches(6.0), Inches(3.0), Inches(0.8), Inches(0.6),
             font_size=Pt(28), color=ACCENT_BLUE, align=PP_ALIGN.CENTER, bold=True)

    add_callout(slide, "每個 Request 都應該可以被任何一台 Server 處理 — 這是 Scale Out 的核心前提",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_12(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "Part 1 小結：架構演進路線圖")
    add_part_label(slide, 1, "傳統部署的演進")
    add_page_number(slide, 12)

    # Timeline baseline
    add_rect(slide, Inches(0.5), Inches(3.15), Inches(12.3), Inches(0.04), ACCENT_BLUE)

    stages = ["單機部署", "DB 分離", "三層架構", "需要 Scale Out"]
    pains = [
        "✅ 快速上線\n⚠ SPOF，無擴展性",
        "✅ 資料持久化\n⚠ 2台機器，網路延遲",
        "✅ 職責清晰\n⚠ 3台機器，部署複雜",
        "✅ 高可用潛力\n⚠ Session / LB 挑戰",
    ]

    node_xs = [0.6, 3.8, 7.0, 10.2]

    for i, (stage, pain) in enumerate(zip(stages, pains)):
        nx = Inches(node_xs[i])
        # Dot
        add_rect(slide, nx, Inches(3.0), Inches(0.32), Inches(0.32), ACCENT_BLUE)
        # Above line: stage title
        add_text(slide, stage, nx - Inches(1.0), Inches(1.85), Inches(2.8), Inches(1.0),
                 font_size=Pt(13), color=TEXT_PRIMARY, align=PP_ALIGN.CENTER, bold=True)
        # Below line: pain points
        add_text(slide, pain, nx - Inches(1.0), Inches(3.35), Inches(2.8), Inches(1.5),
                 font_size=Pt(11), color=TEXT_DIM, align=PP_ALIGN.CENTER)

    add_callout(slide, "下一步 → Part 2：如何 Scale Out，以及隨之而來的新挑戰",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_13(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "Load Balancer：流量分發的核心元件")
    add_part_label(slide, 2, "Scale Out 的挑戰")
    add_page_number(slide, 13)

    # Left: Architecture flow diagram
    add_card(slide, Inches(0.4), Inches(1.6), Inches(5.0), Inches(0.6),
             title=None, body_lines=["Clients (Users)"], accent=TEXT_DIM)
    add_text(slide, "↓", Inches(0.4), Inches(2.25), Inches(5.0), Inches(0.4),
             font_size=Pt(20), color=TEXT_DIM, align=PP_ALIGN.CENTER)
    add_card(slide, Inches(0.4), Inches(2.55), Inches(5.0), Inches(0.7),
             title=None, body_lines=["Load Balancer  :80 / :443"], accent=PART_COLORS[2])
    add_text(slide, "↙  ↓  ↘", Inches(0.4), Inches(3.3), Inches(5.0), Inches(0.4),
             font_size=Pt(18), color=TEXT_DIM, align=PP_ALIGN.CENTER)
    # Three small server cards side by side
    for i, label in enumerate(["App Server 1", "App Server 2", "App Server 3"]):
        add_card(slide, Inches(0.4 + i * 1.6), Inches(3.75), Inches(1.5), Inches(0.55),
                 title=None, body_lines=[label], accent=ACCENT_BLUE)
    add_text(slide, "↓", Inches(0.4), Inches(4.35), Inches(5.0), Inches(0.35),
             font_size=Pt(18), color=TEXT_DIM, align=PP_ALIGN.CENTER)
    add_card(slide, Inches(0.4), Inches(4.65), Inches(5.0), Inches(0.6),
             title=None, body_lines=["Database  :5432"], accent=PART_COLORS[4])

    # Right: Algorithm cards
    algo_cards = [
        (ACCENT_BLUE,         "Round Robin",          ["依序輪流 A→B→C→A  |  適用：同規格 Server"]),
        (PART_COLORS[2],      "Weighted Round Robin", ["依權重比例分配  |  適用：Server 規格不同"]),
        (ACCENT_GREEN,        "Least Connections",    ["優先分配給連線最少的  |  適用：長連線 API"]),
        (ACCENT_AMBER,        "IP Hash",              ["同一 IP 固定導向同台  |  適用：需 Sticky Session"]),
        (RGBColor(0x00,0xE5,0xFF), "Health Check",   ["定期 Ping，自動移除故障節點  |  所有生產環境基本配置"]),
    ]
    for idx, (accent, title, body) in enumerate(algo_cards):
        add_card(slide, Inches(6.0), Inches(1.6 + idx * 0.95), Inches(7.0), Inches(0.85),
                 title=title, body_lines=body, accent=accent)

    add_callout(slide, "LB 是 Scale Out 的入口，健康檢查讓系統具備自動容錯能力",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_14(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "Scale Out 的陷阱：Session 狀態問題")
    add_part_label(slide, 2, "Scale Out 的挑戰")
    add_page_number(slide, 14)

    # Problem card at top
    add_card(slide, Inches(0.4), Inches(1.6), Inches(12.5), Inches(1.5),
             title="問題情境：用戶登入後，下一個 Request 被導到不同的 Server...",
             body_lines=[
                 "Request #1 登入 → Server 1 ✅  |  Request #2 查詢 → Server 2 ❌  |  Request #3 購買 → Server 3 ❌",
                 "結果：Session 散在各自記憶體 → 用戶被強制登出！",
             ],
             accent=RGBColor(0xFF, 0x4A, 0x4A))

    # Three solution cards
    solution_cards = [
        (Inches(0.4),  ACCENT_AMBER, "① Cookie-Based Session",
         ["Session 資料存在 Cookie 中", "每次 Request 帶著資料",
          "✅ 不依賴 Server 狀態", "⚠ 大小有限制 (4KB)", "⚠ 敏感資料需加密"]),
        (Inches(4.6),  ACCENT_GREEN, "② Redis Session Store",
         ["Session 存入 Redis", "所有 Server 共用同一份",
          "✅ 完整 Session 功能", "✅ TTL 自動過期", "⚠ 多一個外部依賴"]),
        (Inches(8.8),  ACCENT_BLUE,  "③ JWT Token",
         ["無狀態 Token，含用戶資訊", "Server 不需存任何東西",
          "✅ 完全 Stateless", "✅ 跨服務驗證", "⚠ 無法立即作廢"]),
    ]
    for x, accent, title, body in solution_cards:
        add_card(slide, x, Inches(3.4), Inches(3.9), Inches(2.6),
                 title=title, body_lines=body, accent=accent)

    add_callout(slide, "核心設計原則：App Server 本身不應儲存任何特定請求的狀態（Stateless）",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_15(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "資料庫的擴展：Read Replica 架構")
    add_part_label(slide, 2, "Scale Out 的挑戰")
    add_page_number(slide, 15)

    # Left half: App server stack
    for i, label in enumerate(["App Server 1", "App Server 2", "App Server 3", "App Server 4"]):
        add_card(slide, Inches(0.4), Inches(1.8 + i * 0.75), Inches(2.8), Inches(0.65),
                 title=None, body_lines=[label], accent=ACCENT_BLUE)

    # Write arrow label
    add_text(slide, "Write →", Inches(3.3), Inches(2.1), Inches(1.2), Inches(0.4),
             font_size=Pt(11), color=ACCENT_AMBER, bold=True)
    # Read arrow label
    add_text(slide, "Read →", Inches(3.3), Inches(3.5), Inches(1.2), Inches(0.4),
             font_size=Pt(11), color=ACCENT_GREEN, bold=True)

    # Center: Primary DB
    add_card(slide, Inches(4.5), Inches(2.3), Inches(3.0), Inches(0.9),
             title="Primary DB (Write Only)", body_lines=[], accent=ACCENT_BLUE)

    # Replication label
    add_text(slide, "Replication →", Inches(4.5), Inches(3.3), Inches(3.0), Inches(0.4),
             font_size=Pt(10), color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # Right: Replica cards
    for i, label in enumerate(["Read Replica 1", "Read Replica 2", "Read Replica 3"]):
        add_card(slide, Inches(7.8), Inches(1.8 + i * 0.75), Inches(3.5), Inches(0.65),
                 title=None, body_lines=[label], accent=TEXT_DIM)

    # Right-side info cards
    add_card(slide, Inches(7.5), Inches(3.8), Inches(5.4), Inches(1.5),
             title="✅ 優點",
             body_lines=["讀取效能大幅提升", "Primary 壓力降低", "容災備援"],
             accent=ACCENT_GREEN)
    add_card(slide, Inches(0.4), Inches(3.8), Inches(6.8), Inches(1.5),
             title="⚠  限制",
             body_lines=["Write 仍是單點", "Replication Lag 資料延遲", "讀寫路由需在 App 層區分"],
             accent=ACCENT_AMBER)

    add_callout(slide, "約 80% 的 DB 操作是讀取。Read Replica 大幅分散壓力。Write 瓶頸需 Sharding 或 NewSQL 解決",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_16(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "三層 Caching 策略：從外到內")
    add_part_label(slide, 2, "Scale Out 的挑戰")
    add_page_number(slide, 16)

    cache_cards = [
        (Inches(0.4),  ACCENT_AMBER,   "① CDN Cache  (最外層)",
         ["靜態資源：JS / CSS / 圖片", "全球節點分發", "TTL: 小時~天", "Cost: 極低",
          "", "命中率: 60~80%", "適合: 不常變動的資源"]),
        (Inches(4.6),  PART_COLORS[2], "② Redis / Memcached  (應用層)",
         ["API 回應結果", "Session Store", "計算結果快取", "TTL: 秒~分鐘",
          "", "命中率: 40~60%", "適合: 熱門查詢結果"]),
        (Inches(8.8),  ACCENT_BLUE,    "③ Local In-Process  (最內層)",
         ["超熱資料 (config / 靜態表)", "microsecond 等級", "TTL: 秒級", "大小: 有限 (幾MB)",
          "", "命中率: 90%+", "適合: 幾乎不變的資料"]),
    ]
    for x, accent, title, body in cache_cards:
        add_card(slide, x, Inches(1.6), Inches(3.9), Inches(3.2),
                 title=title, body_lines=body, accent=accent)

    # Bottom: Cache Invalidation card
    add_card(slide, Inches(0.4), Inches(5.05), Inches(12.5), Inches(1.0),
             title="Cache Invalidation 挑戰",
             body_lines=["快取何時更新？Write-Through / Write-Behind / Cache-Aside 三種策略各有取捨"],
             accent=ACCENT_AMBER)

    add_callout(slide, "Cache 帶來效能，也帶來資料一致性挑戰。There are only two hard things in CS...",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="warning")
    return slide


def slide_17(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "訊息佇列：非同步解耦的利器")
    add_part_label(slide, 2, "Scale Out 的挑戰")
    add_page_number(slide, 17)

    # Left: Sync problem card
    add_card(slide, Inches(0.4), Inches(1.6), Inches(5.8), Inches(2.5),
             title="❌ 同步呼叫的問題",
             body_lines=[
                 "Service A → Service B → Service C",
                 "任一服務慢 → 全部等待",
                 "任一服務掛 → 整條鏈斷",
                 "流量大時 → 下游被壓垮",
                 "Scale 困難：必須一起 Scale",
             ],
             accent=RGBColor(0xFF, 0x4A, 0x4A))

    # Center arrow
    add_text(slide, "→", Inches(6.2), Inches(2.6), Inches(0.8), Inches(0.6),
             font_size=Pt(32), color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    # Right: MQ solution card
    add_card(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(2.5),
             title="✅ MQ 解耦",
             body_lines=[
                 "Producer → Queue → Consumer",
                 "Producer 不知道 Consumer 是誰",
                 "Consumer 可以獨立 Scale",
                 "Consumer 掛掉 → 訊息在 Queue 等待",
                 "解耦 = 獨立部署、獨立容錯",
             ],
             accent=ACCENT_GREEN)

    # Bottom use case cards
    use_cases = [
        (Inches(0.4),  "Email / 通知發送",   ["非同步處理，不阻塞主流程"]),
        (Inches(4.6),  "影片 / 圖片處理",    ["耗時任務排隊，Consumer 依序處理"]),
        (Inches(8.8),  "訂單 / 金流流程",    ["確保不遺失，重試機制"]),
    ]
    for x, title, body in use_cases:
        add_card(slide, x, Inches(4.4), Inches(3.9), Inches(1.6),
                 title=title, body_lines=body, accent=ACCENT_BLUE)

    add_callout(slide, "MQ 讓服務彼此獨立，任一方可以獨立 Scale 且不互相影響",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_18(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "完整分散式架構全貌")
    add_part_label(slide, 2, "Scale Out 的挑戰")
    add_page_number(slide, 18)

    # Row 1
    add_card(slide, Inches(0.4), Inches(1.5), Inches(12.5), Inches(0.5),
             title=None,
             body_lines=["Internet Users → CDN (Cloudflare) → WAF / Firewall → DNS LB"],
             accent=TEXT_DIM)
    # Arrow
    add_text(slide, "↓", Inches(6.0), Inches(2.05), Inches(1.5), Inches(0.25),
             font_size=Pt(14), color=TEXT_DIM, align=PP_ALIGN.CENTER)
    # Row 2
    add_card(slide, Inches(4.5), Inches(2.15), Inches(4.0), Inches(0.5),
             title=None, body_lines=["Frontend LB (Nginx)"], accent=PART_COLORS[3])
    # Arrow
    add_text(slide, "↓", Inches(6.0), Inches(2.7), Inches(1.5), Inches(0.25),
             font_size=Pt(14), color=TEXT_DIM, align=PP_ALIGN.CENTER)
    # Row 3: 3 frontend servers
    for i, label in enumerate(["Frontend Server 1", "Frontend Server 2", "Frontend Server 3"]):
        add_card(slide, Inches(0.3 + i * 4.2), Inches(2.8), Inches(3.8), Inches(0.5),
                 title=None, body_lines=[label], accent=PART_COLORS[3])
    # Arrow
    add_text(slide, "↓", Inches(6.0), Inches(3.35), Inches(1.5), Inches(0.25),
             font_size=Pt(14), color=TEXT_DIM, align=PP_ALIGN.CENTER)
    # Row 4
    add_card(slide, Inches(4.5), Inches(3.45), Inches(4.0), Inches(0.5),
             title=None, body_lines=["Backend LB (HAProxy)"], accent=PART_COLORS[2])
    # Arrow
    add_text(slide, "↓", Inches(6.0), Inches(4.0), Inches(1.5), Inches(0.25),
             font_size=Pt(14), color=TEXT_DIM, align=PP_ALIGN.CENTER)
    # Row 5: 3 backend servers
    for i, label in enumerate(["Backend Server 1", "Backend Server 2", "Backend Server 3"]):
        add_card(slide, Inches(0.3 + i * 4.2), Inches(4.1), Inches(3.8), Inches(0.5),
                 title=None, body_lines=[label], accent=PART_COLORS[2])
    # Arrow
    add_text(slide, "↓", Inches(6.0), Inches(4.65), Inches(1.5), Inches(0.25),
             font_size=Pt(14), color=TEXT_DIM, align=PP_ALIGN.CENTER)
    # Row 6: infra
    add_card(slide, Inches(0.3),  Inches(4.75), Inches(3.8), Inches(0.5),
             title=None, body_lines=["Redis Cluster"], accent=ACCENT_AMBER)
    add_card(slide, Inches(4.5),  Inches(4.75), Inches(3.8), Inches(0.5),
             title=None, body_lines=["MQ (RabbitMQ)"], accent=ACCENT_AMBER)
    add_card(slide, Inches(8.7),  Inches(4.75), Inches(4.2), Inches(0.5),
             title=None, body_lines=["DB Primary + 3 Replica"], accent=PART_COLORS[4])

    add_callout(slide,
                "共需維護：2 LB + 3 Frontend + 3 Backend + 4 Cache/MQ + 4 DB + CDN/WAF = 15+ 台機器！",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="danger")
    return slide


def slide_19(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "部署惡夢：維運複雜度爆炸")
    add_part_label(slide, 2, "Scale Out 的挑戰")
    add_page_number(slide, 19)

    danger_color = RGBColor(0xFF, 0x4A, 0x4A)
    grid_cards = [
        (Inches(0.4), Inches(1.6),  "⑴ 環境不一致",
         ["Dev：「我機器上能跑」", "15 台機器的 Runtime 版本可能都不同", "手動 SSH 逐一更新，出錯機率極高"]),
        (Inches(6.9), Inches(1.6),  "⑵ 更新 / 回滾困難",
         ["新版本需依序更新全部 15 台機器", "某台失敗了怎麼回滾？", "停機時間難以控制，風險極高"]),
        (Inches(0.4), Inches(3.8),  "⑶ 設定管理混亂",
         ["每台機器的 config 可能不同", "DB 連線字串、Secret Key 散落各處", "新機器如何確保設定正確？"]),
        (Inches(6.9), Inches(3.8),  "⑷ Dev/Prod 環境落差",
         ["開發者用 macOS，Server 是 Ubuntu", "套件相依性（Dependency Hell）", "「在本地跑得好好的...」"]),
    ]
    for x, y, title, body in grid_cards:
        add_card(slide, x, y, Inches(5.8), Inches(2.0),
                 title=title, body_lines=body, accent=danger_color)

    add_callout(slide, "資深工程師每天花大量時間在「機器設定與維護」而非「產品開發」",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="warning")
    return slide


def slide_20(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "Part 2 小結：問題累積，需要新思維")
    add_part_label(slide, 2, "Scale Out 的挑戰")
    add_page_number(slide, 20)

    summary_rows = [
        (ACCENT_BLUE,        "Load Balancer  ✅ 單一入口、健康檢查  ⚠ 需 Stateless 設計、新增設定層"),
        (PART_COLORS[2],     "Session Store  ✅ 解決多機 Session  ⚠ 多一個 Redis 依賴要維護"),
        (ACCENT_GREEN,       "DB Read Replica  ✅ 讀取效能大幅提升  ⚠ Write 仍是瓶頸、Replication Lag"),
        (ACCENT_AMBER,       "Caching  ✅ 大幅降低 DB 壓力  ⚠ Cache 一致性問題、過期策略複雜"),
        (PART_COLORS[3],     "訊息佇列 MQ  ✅ 服務解耦、非同步  ⚠ 增加基礎設施複雜度"),
        (PART_COLORS[4],     "完整分散式架構  ✅ 高可用 + 高效能  ⚠ 15+ 台機器，維運地獄"),
    ]
    for idx, (accent, text) in enumerate(summary_rows):
        add_card(slide, Inches(0.4), Inches(1.6 + idx * 0.7), Inches(12.5), Inches(0.62),
                 title=None, body_lines=[text], accent=accent)

    add_callout(slide, "有沒有辦法，讓「在哪台機器上跑」不再是問題？→ 答案就是：Container！",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_21(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "Container 是什麼？解決什麼問題？")
    add_part_label(slide, 3, "Container 革命")
    add_page_number(slide, 21)

    # Left half: traditional deployment problem
    add_card(slide, Inches(0.4), Inches(1.6), Inches(5.8), Inches(3.5),
             title="❌ 傳統部署的困境",
             body_lines=[
                 "應用程式 Code",
                 "Runtime  Python 3.9",
                 "系統套件  apt packages",
                 "OS Config  .conf files",
                 "Ubuntu 20.04 / 機器 A",
                 "",
                 "換台機器 = 環境可能完全不同！",
             ],
             accent=RGBColor(0xFF, 0x4A, 0x4A))

    # Center arrow
    add_text(slide, "→", Inches(6.3), Inches(3.0), Inches(0.8), Inches(0.6),
             font_size=Pt(32), color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    # Right half: container solution
    add_card(slide, Inches(7.2), Inches(1.6), Inches(5.8), Inches(3.5),
             title="✅ Container 解法",
             body_lines=[
                 "Container Image（封裝一切）",
                 "  ├ 應用程式 Code",
                 "  ├ Runtime + 所有套件",
                 "  └ OS 最小依賴",
                 "",
                 "Container Runtime (Docker)",
                 "任何 Linux 主機",
                 "",
                 "Build once, run anywhere!",
             ],
             accent=ACCENT_GREEN)

    # Bottom center: three core values card
    add_card(slide, Inches(0.4), Inches(5.35), Inches(12.5), Inches(1.0),
             title="三大核心價值",
             body_lines=["① 環境一致性：開發 = 測試 = 生產  ② 快速啟動：秒級甚至毫秒級  ③ 輕量隔離：不需要完整的 OS"],
             accent=PART_COLORS[3])

    # Callout tip
    add_callout(slide, "Container 不是 VM！更輕量、更快速、更適合現代雲端環境",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_22(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "VM vs Container：架構深度對比")
    add_part_label(slide, 3, "Container 革命")
    add_page_number(slide, 22)

    # Left half: VM
    add_text(slide, "Virtual Machine (VM)", Inches(0.4), Inches(1.55), Inches(5.8), Inches(0.4),
             font_name=FONT_TITLE, font_size=Pt(16), color=ACCENT_AMBER, bold=True)

    # VM layered boxes
    add_card(slide, Inches(0.4), Inches(2.0), Inches(2.7), Inches(0.6),
             body_lines=["App A"], accent=ACCENT_BLUE)
    add_card(slide, Inches(3.3), Inches(2.0), Inches(2.7), Inches(0.6),
             body_lines=["App B"], accent=ACCENT_BLUE)
    add_card(slide, Inches(0.4), Inches(2.7), Inches(2.7), Inches(0.6),
             body_lines=["Guest OS A (Ubuntu 20)"], accent=CARD_BORDER)
    add_card(slide, Inches(3.3), Inches(2.7), Inches(2.7), Inches(0.6),
             body_lines=["Guest OS B (CentOS 8)"], accent=CARD_BORDER)
    add_card(slide, Inches(0.4), Inches(3.4), Inches(5.8), Inches(0.6),
             body_lines=["Hypervisor (VMware / KVM)"], accent=ACCENT_AMBER)
    add_card(slide, Inches(0.4), Inches(4.1), Inches(5.8), Inches(0.6),
             body_lines=["Host OS (Linux / Windows)"], accent=TEXT_DIM)
    add_card(slide, Inches(0.4), Inches(4.75), Inches(5.8), Inches(0.6),
             body_lines=["Physical Hardware / Cloud"], accent=CARD_BORDER)

    # VM stats
    vm_stats = [
        "⏱ 啟動時間：分鐘級（需 boot OS）",
        "💾 Image 大小：GB 級（含完整 Guest OS）",
        "🔒 隔離強度：強（完整 OS 層隔離）",
    ]
    for i, stat in enumerate(vm_stats):
        add_text(slide, stat, Inches(0.4), Inches(5.5 + i * 0.28), Inches(5.8), Inches(0.28),
                 font_size=Pt(10), color=TEXT_DIM)

    # Right half: Container
    add_text(slide, "Container", Inches(7.0), Inches(1.55), Inches(5.8), Inches(0.4),
             font_name=FONT_TITLE, font_size=Pt(16), color=ACCENT_GREEN, bold=True)

    # Container layered boxes
    add_card(slide, Inches(7.0), Inches(2.0), Inches(2.7), Inches(0.6),
             body_lines=["Container A (App A + Libs)"], accent=PART_COLORS[3])
    add_card(slide, Inches(9.8), Inches(2.0), Inches(2.7), Inches(0.6),
             body_lines=["Container B (App B + Libs)"], accent=PART_COLORS[3])
    add_card(slide, Inches(7.0), Inches(2.7), Inches(5.8), Inches(0.6),
             body_lines=["Container Runtime (Docker / containerd)"], accent=ACCENT_GREEN)
    add_card(slide, Inches(7.0), Inches(3.4), Inches(5.8), Inches(0.6),
             body_lines=["Host OS Kernel（共用，不重複安裝）"], accent=TEXT_DIM)
    add_card(slide, Inches(7.0), Inches(4.1), Inches(5.8), Inches(0.6),
             body_lines=["Physical Hardware / Cloud VM"], accent=CARD_BORDER)

    # Container stats
    ct_stats = [
        "⚡ 啟動時間：秒級甚至毫秒級",
        "📦 Image 大小：MB 級（只含必要檔案）",
        "🔐 隔離強度：程序級（Namespace + cgroups）",
    ]
    for i, stat in enumerate(ct_stats):
        add_text(slide, stat, Inches(7.0), Inches(5.5 + i * 0.28), Inches(5.8), Inches(0.28),
                 font_size=Pt(10), color=TEXT_DIM)

    # Center VS divider
    add_rect(slide, Inches(6.45), Inches(1.6), Inches(0.04), Inches(4.8), CARD_BORDER)

    # Callout tip
    add_callout(slide, "Container 共用 Host OS Kernel，更輕量快速。安全需求極高時才考慮 VM（或兩者並用）",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_23(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "Docker 核心概念：Image → Container")
    add_part_label(slide, 3, "Container 革命")
    add_page_number(slide, 23)

    # Top flow: Dockerfile
    add_card(slide, Inches(0.4), Inches(1.6), Inches(2.5), Inches(1.0),
             body_lines=["FROM python:3.11-slim", "COPY . /app", "RUN pip install -r req.txt", "CMD [\"python\",\"main.py\"]"],
             accent=TEXT_DIM)
    add_text(slide, "docker build →", Inches(3.1), Inches(1.85), Inches(1.6), Inches(0.4),
             font_size=Pt(12), color=TEXT_DIM)
    add_card(slide, Inches(4.8), Inches(1.6), Inches(3.0), Inches(1.0),
             body_lines=["Docker Image（不可變快照）"],
             accent=PART_COLORS[3])
    add_text(slide, "docker run →", Inches(8.0), Inches(1.85), Inches(1.1), Inches(0.4),
             font_size=Pt(12), color=TEXT_DIM)
    add_card(slide, Inches(9.2), Inches(1.6), Inches(3.7), Inches(1.0),
             body_lines=["Container（執行中的程序）"],
             accent=ACCENT_GREEN)

    # Below Image card: push/pull + registry
    add_text(slide, "push ↑ / pull ↓", Inches(4.8), Inches(2.65), Inches(3.0), Inches(0.35),
             font_size=Pt(11), color=TEXT_DIM, align=PP_ALIGN.CENTER)
    add_card(slide, Inches(4.8), Inches(3.0), Inches(3.0), Inches(0.8),
             body_lines=["Container Registry", "Docker Hub / ECR"],
             accent=ACCENT_AMBER)

    # Bottom: common docker commands card
    add_card(slide, Inches(0.4), Inches(3.4), Inches(12.5), Inches(3.5),
             title="常用 Docker 指令",
             body_lines=[
                 "docker build -t my-app:v1 .    → 從 Dockerfile 建立 Image",
                 "docker run -d -p 8080:8080 my-app:v1    → 背景執行，映射 Port",
                 "docker ps    → 列出執行中的 Container",
                 "docker logs <id>    → 查看 Container 日誌",
                 "docker exec -it <id> /bin/bash    → 進入 Container Shell",
                 "docker stop <id> && docker rm <id>    → 停止並刪除",
             ],
             accent=ACCENT_BLUE)
    return slide


def slide_24(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "Docker Compose：多容器應用一鍵啟動")
    add_part_label(slide, 3, "Container 革命")
    add_page_number(slide, 24)

    # Left: docker-compose.yml card
    add_card(slide, Inches(0.4), Inches(1.6), Inches(6.0), Inches(4.8),
             title="docker-compose.yml",
             body_lines=[
                 "version: '3.8'",
                 "services:",
                 "  frontend:",
                 "    image: nginx:alpine",
                 "    ports: [\"80:80\"]",
                 "  backend:",
                 "    build: ./backend",
                 "    ports: [\"8080:8080\"]",
                 "    environment:",
                 "      DB_HOST: db",
                 "      SECRET_KEY: ${SECRET}",
                 "    depends_on: [db, redis]",
                 "  db:",
                 "    image: postgres:15",
                 "    volumes:",
                 "      - pgdata:/var/lib/pg",
                 "  redis:",
                 "    image: redis:7-alpine",
                 "volumes:",
                 "  pgdata:",
             ],
             accent=PART_COLORS[3])

    # Right: command card
    add_card(slide, Inches(6.8), Inches(1.6), Inches(6.0), Inches(0.8),
             body_lines=["$ docker compose up -d"],
             accent=ACCENT_GREEN)

    # Label
    add_text(slide, "一個指令，啟動完整應用：", Inches(6.8), Inches(2.55), Inches(6.0), Inches(0.28),
             font_size=Pt(12), color=TEXT_DIM)

    # 4 service status cards
    service_cards = [
        (PART_COLORS[3], "frontend  (Nginx : 80)"),
        (ACCENT_BLUE,    "backend   (Python : 8080)"),
        (PART_COLORS[4], "db        (PostgreSQL : 5432)"),
        (ACCENT_AMBER,   "redis     (Redis : 6379)"),
    ]
    for i, (accent, text) in enumerate(service_cards):
        add_card(slide, Inches(6.8), Inches(2.75 + i * 0.72), Inches(5.8), Inches(0.65),
                 body_lines=[text], accent=accent)

    # Volume note
    add_card(slide, Inches(6.8), Inches(5.65), Inches(5.8), Inches(0.6),
             body_lines=["Volume: pgdata — 資料持久化，容器重建不遺失"],
             accent=TEXT_DIM)

    # Callout
    add_callout(slide, "所有容器在同一虛擬網路，用服務名稱互相連線（如 DB_HOST: db），不需寫 IP",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_25(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "Container Registry：Image 的倉庫與版本管理")
    add_part_label(slide, 3, "Container 革命")
    add_page_number(slide, 25)

    # Top flow: build pipeline
    # Cards
    add_card(slide, Inches(0.4), Inches(1.65), Inches(1.8), Inches(0.65),
             body_lines=["Code + Dockerfile"], accent=TEXT_DIM)
    add_text(slide, "→", Inches(2.25), Inches(1.7), Inches(0.6), Inches(0.55),
             font_size=Pt(20), color=TEXT_DIM, align=PP_ALIGN.CENTER)
    add_card(slide, Inches(3.5), Inches(1.65), Inches(1.8), Inches(0.65),
             body_lines=["Image"], accent=PART_COLORS[3])
    add_text(slide, "→", Inches(5.25), Inches(1.7), Inches(0.6), Inches(0.55),
             font_size=Pt(20), color=TEXT_DIM, align=PP_ALIGN.CENTER)
    add_card(slide, Inches(6.4), Inches(1.65), Inches(1.8), Inches(0.65),
             body_lines=["Registry"], accent=ACCENT_AMBER)
    add_text(slide, "→", Inches(8.15), Inches(1.7), Inches(0.6), Inches(0.55),
             font_size=Pt(20), color=TEXT_DIM, align=PP_ALIGN.CENTER)
    add_card(slide, Inches(11.2), Inches(1.65), Inches(1.8), Inches(0.65),
             body_lines=["Server / K8s"], accent=ACCENT_GREEN)

    # Labels between cards
    add_text(slide, "docker build", Inches(2.3), Inches(2.3), Inches(1.1), Inches(0.3),
             font_size=Pt(10), color=TEXT_DIM, align=PP_ALIGN.CENTER)
    add_text(slide, "docker push", Inches(5.3), Inches(2.3), Inches(1.1), Inches(0.3),
             font_size=Pt(10), color=TEXT_DIM, align=PP_ALIGN.CENTER)
    add_text(slide, "docker pull", Inches(8.2), Inches(2.3), Inches(1.1), Inches(0.3),
             font_size=Pt(10), color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # Middle: two cards side by side
    add_card(slide, Inches(0.4), Inches(3.0), Inches(5.8), Inches(1.8),
             title="Registry 選項",
             body_lines=[
                 "Docker Hub — 公開免費，私有收費",
                 "AWS ECR — 與 AWS 整合",
                 "GCR (Google) / ACR (Azure)",
                 "Harbor — 自建私有 Registry",
             ],
             accent=ACCENT_BLUE)
    add_card(slide, Inches(6.6), Inches(3.0), Inches(5.8), Inches(1.8),
             title="Image Tag 策略",
             body_lines=[
                 ":latest — ⚠ 生產環境禁用！",
                 ":v1.2.3 — 語義化版本，推薦",
                 ":git-abc1234 — Git commit SHA",
                 ":env-prod / :env-staging — 環境標籤",
             ],
             accent=PART_COLORS[2])

    # Callout warning
    add_callout(slide, "永遠不要在生產環境用 :latest tag！無法追溯版本、無法確定性回滾",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="warning")
    return slide


def slide_26(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "Container 化：從維運噩夢到一致部署")
    add_part_label(slide, 3, "Container 革命")
    add_page_number(slide, 26)

    # Left header
    add_text(slide, "❌ 傳統：15+ 台機器，各自設定",
             Inches(0.4), Inches(1.55), Inches(5.8), Inches(0.35),
             font_size=Pt(14), color=RGBColor(0xFF, 0x4A, 0x4A), bold=True)

    # 6 small cards stacked
    old_servers = [
        "Server A  Python 3.8 / Ubuntu 20",
        "Server B  Python 3.9 / Ubuntu 22",
        "Server C  Python 3.8 / CentOS 8",
        "DB Server  PostgreSQL 13",
        "Cache  Redis 6",
        "LB  Nginx 1.18",
    ]
    for i, text in enumerate(old_servers):
        add_card(slide, Inches(0.4), Inches(1.9 + i * 0.65), Inches(5.8), Inches(0.58),
                 body_lines=[text], accent=CARD_BORDER)

    # Problem text
    add_text(slide, "問題：每台 SSH 進去設定，每次更新都是手動操作的噩夢",
             Inches(0.4), Inches(5.9), Inches(5.8), Inches(0.4),
             font_size=Pt(11), color=RGBColor(0xFF, 0x4A, 0x4A))

    # Center arrow
    add_text(slide, "→", Inches(6.3), Inches(3.3), Inches(0.8), Inches(0.7),
             font_size=Pt(36), color=ACCENT_BLUE, align=PP_ALIGN.CENTER, bold=True)

    # Right header
    add_text(slide, "✅ Container：統一封裝，一致部署",
             Inches(7.0), Inches(1.55), Inches(5.8), Inches(0.35),
             font_size=Pt(14), color=ACCENT_GREEN, bold=True)

    # 4 service image cards
    new_services = [
        (PART_COLORS[3], "frontend:v2.1  (nginx:alpine)"),
        (ACCENT_BLUE,    "backend:v2.1  (python:3.11-slim)"),
        (PART_COLORS[4], "db:latest  (postgres:15)"),
        (ACCENT_AMBER,   "redis:7  (redis:alpine)"),
    ]
    for i, (accent, text) in enumerate(new_services):
        add_card(slide, Inches(7.0), Inches(1.9 + i * 0.65), Inches(5.8), Inches(0.58),
                 body_lines=[text], accent=accent)

    # Solution card
    add_card(slide, Inches(7.0), Inches(4.7), Inches(5.8), Inches(1.5),
             body_lines=[
                 "一個 Dockerfile 定義環境",
                 "docker compose up 一鍵啟動",
                 "環境完全一致：開發 = 測試 = 生產",
             ],
             accent=ACCENT_GREEN)

    # Callout tip
    add_callout(slide, "Container 解決的是「環境一致性」問題。下一步：如何讓應用程式天生適合 Container 化？",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_27(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "什麼是 Cloud-Ready 的應用？")
    add_part_label(slide, 4, "12-Factor App")
    add_page_number(slide, 27)

    # Left card: ❌ 能跑 ≠ 能 Scale
    add_card(slide, Inches(0.4), Inches(1.6), Inches(5.8), Inches(3.0),
             title="❌ 能跑 ≠ 能 Scale",
             body_lines=[
                 "App 能在一台機器跑",
                 "≠",
                 "App 能在 10 台機器正確運行",
                 "",
                 "常見問題：",
                 "• Config 寫死在 Code 裡",
                 "• Log 只寫到本地檔案",
                 "• Session 存在記憶體",
             ],
             accent=RGBColor(0xFF, 0x4A, 0x4A))

    # Right card: ✅ 12-Factor App
    add_card(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(3.0),
             title="✅ 12-Factor App",
             body_lines=[
                 "2011 年由 Heroku 工程師 Adam Wiggins 提出",
                 "",
                 "12 條設計原則，讓應用：",
                 "• 天生適合 Container 化",
                 "• 易於水平 Scale",
                 "• 環境一致，可預期部署",
                 "• 易於維護和觀測",
             ],
             accent=ACCENT_GREEN)

    # Center connector arrow
    add_text(slide, "→ 設計哲學", Inches(6.0), Inches(2.8), Inches(0.8), Inches(0.4),
             font_size=Pt(12), color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # Bottom card
    add_card(slide, Inches(0.4), Inches(4.85), Inches(12.5), Inches(1.3),
             title="為什麼現在學 12-Factor？",
             body_lines=[
                 "Container 解決「在哪跑」的問題  →  12-Factor 解決「怎麼寫才適合雲端」的問題",
                 "兩者相輔相成：Container 是容器，12-Factor 是設計良好的內容物",
             ],
             accent=PART_COLORS[4])

    # Callout tip
    add_callout(slide, "12-Factor 不是技術規範，是設計哲學。違反它的代價，會在 Scale 時加倍奉還",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_28(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "12-Factor App 全覽")
    add_part_label(slide, 4, "12-Factor App")
    add_page_number(slide, 28)

    factors = [
        ("I. Codebase",          "一個 Repo，多個部署環境",          ACCENT_BLUE),
        ("II. Dependencies",     "明確宣告所有相依套件",              ACCENT_BLUE),
        ("III. Config",          "設定存環境變數，不 hardcode",       ACCENT_BLUE),
        ("IV. Backing Services", "DB / Redis 視為可替換外部資源",     PART_COLORS[4]),
        ("V. Build-Release-Run", "三階段嚴格分離",                    PART_COLORS[4]),
        ("VI. Processes",        "無狀態 Process，狀態外部化",        PART_COLORS[4]),
        ("VII. Port Binding",    "App 自帶 HTTP Server",              PART_COLORS[2]),
        ("VIII. Concurrency",    "以 Process 模型水平擴展",           PART_COLORS[2]),
        ("IX. Disposability",    "快速啟動，優雅關閉",                PART_COLORS[2]),
        ("X. Dev/Prod Parity",   "三環境盡量一致",                    ACCENT_GREEN),
        ("XI. Logs",             "Log 輸出到 stdout/stderr",          ACCENT_GREEN),
        ("XII. Admin Processes", "管理任務作為一次性 Process",        ACCENT_GREEN),
    ]

    card_w = Inches(4.0)
    card_h = Inches(1.2)
    x_starts = [Inches(0.3), Inches(4.5), Inches(8.7)]
    y_start = Inches(1.55)
    y_spacing = Inches(1.3)

    for i, (factor_name, desc, accent) in enumerate(factors):
        row = i // 3
        col = i % 3
        x = x_starts[col]
        y = y_start + row * y_spacing
        add_card(slide, x, y, card_w, card_h,
                 title=factor_name,
                 body_lines=[desc],
                 accent=accent)

    return slide


def slide_29(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "Factor 1–3：Codebase / Dependencies / Config")
    add_part_label(slide, 4, "12-Factor App")
    add_page_number(slide, 29)

    card_w = Inches(4.0)
    card_h = Inches(4.5)
    x_positions = [Inches(0.3), Inches(4.6), Inches(8.9)]

    cards = [
        ("I. Codebase", ACCENT_BLUE, [
            "原則：一個 App = 一個 Repo",
            "",
            "✅ 正確：",
            "  一個 Repo，多個部署",
            "  (dev / staging / prod)",
            "",
            "❌ 錯誤：",
            "  多個 Repo 共享同一份 Code",
            "",
            "→ 多個 App 用 Monorepo 時",
            "   每個 App 仍需獨立部署單元",
        ]),
        ("II. Dependencies", ACCENT_BLUE, [
            "原則：明確宣告，隔離安裝",
            "",
            "✅ 正確：",
            "  requirements.txt / pyproject.toml",
            "  package.json + package-lock.json",
            "",
            "❌ 錯誤：",
            "  依賴系統已安裝的套件",
            "  (「反正 Server 上有 curl」)",
            "",
            "→ Container 天然支援此 Factor",
        ]),
        ("III. Config", ACCENT_BLUE, [
            "原則：設定存環境變數",
            "",
            "✅ 正確：",
            "  DB_HOST=localhost",
            "  SECRET_KEY=$SECRET",
            "",
            "❌ 錯誤：",
            "  DB_HOST = 'db.prod.internal'",
            "  寫死在 config.py 裡",
            "",
            "→ 不同環境只改 env var，",
            "   不改 Code 或 Image",
        ]),
    ]

    for i, (title, accent, body_lines) in enumerate(cards):
        add_card(slide, x_positions[i], Inches(1.55), card_w, card_h,
                 title=title, body_lines=body_lines, accent=accent)

    return slide


def slide_30(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "Factor 4–6：Backing Services / Build-Release-Run / Processes")
    add_part_label(slide, 4, "12-Factor App")
    add_page_number(slide, 30)

    card_w = Inches(4.0)
    card_h = Inches(4.5)
    x_positions = [Inches(0.3), Inches(4.6), Inches(8.9)]

    cards = [
        ("IV. Backing Services", PART_COLORS[4], [
            "原則：外部資源視為可替換",
            "",
            "Backing Services：",
            "  PostgreSQL / MySQL",
            "  Redis / Memcached",
            "  S3 / 任何外部 API",
            "",
            "✅ 正確：",
            "  切換 DB 只改 DB_URL",
            "  不需修改 App Code",
            "",
            "→ 讓測試環境換成 SQLite",
            "   生產換成 PostgreSQL 成為可能",
        ]),
        ("V. Build-Release-Run", PART_COLORS[4], [
            "原則：三階段嚴格分離",
            "",
            "Build：Code → 可執行 Artifact",
            "  (docker build)",
            "",
            "Release：Artifact + Config",
            "  (Image + env vars)",
            "",
            "Run：執行 Release",
            "  (docker run / k8s deploy)",
            "",
            "❌ 錯誤：",
            "  上線後直接 ssh 進去改 Code",
        ]),
        ("VI. Processes", PART_COLORS[4], [
            "原則：無狀態 Process",
            "",
            "✅ 正確：",
            "  App 是無狀態的",
            "  Session → Redis",
            "  上傳檔案 → S3",
            "",
            "❌ 錯誤：",
            "  檔案存在本地 /tmp",
            "  計算結果存在記憶體",
            "",
            "→ 任何一台 Server 掛掉",
            "   不影響其他 Server 的服務",
        ]),
    ]

    for i, (title, accent, body_lines) in enumerate(cards):
        add_card(slide, x_positions[i], Inches(1.55), card_w, card_h,
                 title=title, body_lines=body_lines, accent=accent)

    return slide


def slide_31(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "Factor 7–9：Port Binding / Concurrency / Disposability")
    add_part_label(slide, 4, "12-Factor App")
    add_page_number(slide, 31)

    card_w = Inches(4.0)
    card_h = Inches(4.5)
    x_positions = [Inches(0.3), Inches(4.6), Inches(8.9)]

    cards = [
        ("VII. Port Binding", PART_COLORS[2], [
            "原則：App 自帶 HTTP Server",
            "",
            "✅ 正確：",
            "  App 直接 listen port",
            "  (FastAPI 內建 uvicorn)",
            "  (Node.js 直接 listen 8080)",
            "",
            "❌ 錯誤：",
            "  依賴外部 Apache / mod_wsgi",
            "",
            "→ App 可以直接被 LB 呼叫",
            "   不需要中間的 Web Server 層",
        ]),
        ("VIII. Concurrency", PART_COLORS[2], [
            "原則：以 Process 模型 Scale",
            "",
            "水平擴展 = 加 Process：",
            "  web: 3 個 Process",
            "  worker: 5 個 Process",
            "",
            "✅ 正確：",
            "  每種工作類型獨立 Scale",
            "",
            "❌ 錯誤：",
            "  單一超級大 Process 做所有事",
            "",
            "→ 對應 Docker Compose / K8s replica",
        ]),
        ("IX. Disposability", PART_COLORS[2], [
            "原則：快速啟動，優雅關閉",
            "",
            "快速啟動：",
            "  啟動時間 < 幾秒",
            "  (Container 冷啟動問題)",
            "",
            "優雅關閉：",
            "  收到 SIGTERM → 完成當前請求",
            "  → 拒絕新請求 → 關閉",
            "",
            "→ K8s 可以隨時殺掉 Pod",
            "   App 必須能處理這種情況",
        ]),
    ]

    for i, (title, accent, body_lines) in enumerate(cards):
        add_card(slide, x_positions[i], Inches(1.55), card_w, card_h,
                 title=title, body_lines=body_lines, accent=accent)

    return slide


def slide_32(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "Factor 10–12：Dev/Prod Parity / Logs / Admin")
    add_part_label(slide, 4, "12-Factor App")
    add_page_number(slide, 32)

    card_w = Inches(4.0)
    card_h = Inches(4.5)
    x_positions = [Inches(0.3), Inches(4.6), Inches(8.9)]

    cards = [
        ("X. Dev/Prod Parity", ACCENT_GREEN, [
            "原則：三環境盡量一致",
            "",
            "時間差：",
            "  Dev → Prod 盡快（CI/CD）",
            "",
            "人員差：",
            "  Dev 也要了解 Prod 狀況",
            "",
            "工具差：",
            "  ❌ 本地 SQLite + 生產 PostgreSQL",
            "  ✅ 本地也用 Docker + PostgreSQL",
            "",
            "→ Container + Docker Compose",
            "   讓本地環境 = 生產環境",
        ]),
        ("XI. Logs", ACCENT_GREEN, [
            "原則：Log 輸出到 stdout",
            "",
            "✅ 正確：",
            "  print() / logger.info()",
            "  App 不管 Log 要去哪裡",
            "",
            "❌ 錯誤：",
            "  Log 寫到 /var/log/app.log",
            "  App 自己管 Log rotation",
            "",
            "→ 容器平台負責收集 Log",
            "   (ELK / Loki / CloudWatch)",
            "",
            "→ App 職責單一，不管 Log 路由",
        ]),
        ("XII. Admin Processes", ACCENT_GREEN, [
            "原則：管理任務一次性執行",
            "",
            "常見管理任務：",
            "  DB Migration",
            "  資料修復腳本",
            "  一次性初始化",
            "",
            "✅ 正確：",
            "  docker run myapp python migrate.py",
            "  k8s Job / CronJob",
            "",
            "❌ 錯誤：",
            "  SSH 進去手動執行",
            "  混在 App 啟動流程裡",
        ]),
    ]

    for i, (title, accent, body_lines) in enumerate(cards):
        add_card(slide, x_positions[i], Inches(1.55), card_w, card_h,
                 title=title, body_lines=body_lines, accent=accent)

    return slide


def slide_33(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "常見違反 12-Factor 的案例分析")
    add_part_label(slide, 4, "12-Factor App")
    add_page_number(slide, 33)

    danger_color = RGBColor(0xFF, 0x4A, 0x4A)

    cases = [
        (Inches(0.4), Inches(1.6), "案例 1：Config 寫死在 Code",
         ["DB_URL = 'postgresql://admin:password@db.prod/'",
          "發現問題：換 DB 要改 Code → 改 Code 要重新 Build → 違反 Factor III & V"]),
        (Inches(6.9), Inches(1.6), "案例 2：Session 存在記憶體",
         ["users = {} # {user_id: session_data}",
          "發現問題：Scale Out 後 Session 消失 → 用戶被強制登出 → 違反 Factor VI"]),
        (Inches(0.4), Inches(3.8), "案例 3：Log 寫到本地檔案",
         ["logging.FileHandler('/var/log/app.log')",
          "發現問題：K8s Pod 重啟 Log 遺失、分散式環境無法集中查看 → 違反 Factor XI"]),
        (Inches(6.9), Inches(3.8), "案例 4：本地 SQLite + 生產 PG",
         ["# dev: sqlite:///./app.db  # prod: postgresql://...",
          "發現問題：SQLite 和 PostgreSQL 行為差異 → 本地測試過，上線就壞 → 違反 Factor X"]),
    ]

    for x, y, title, body in cases:
        add_card(slide, x, y, Inches(5.8), Inches(2.0),
                 title=title, body_lines=body, accent=danger_color)

    add_callout(slide, "每一個違反，都是未來在不對的時間點爆發的技術債",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="warning")
    return slide


def slide_34(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "12-Factor 與分散式架構的對應關係")
    add_part_label(slide, 4, "12-Factor App")
    add_page_number(slide, 34)

    left_items = [
        ("Load Balancer",               PART_COLORS[2]),
        ("App Server (Stateless)",      ACCENT_BLUE),
        ("Config / Secret 管理",        ACCENT_BLUE),
        ("Session Store (Redis)",       ACCENT_AMBER),
        ("Log 聚合平台",                ACCENT_GREEN),
        ("Container Registry",          PART_COLORS[3]),
    ]

    right_items = [
        ("→ Factor IX：Disposability 快速啟動/關閉", PART_COLORS[2]),
        ("→ Factor VI：Processes 無狀態",            ACCENT_BLUE),
        ("→ Factor III：Config 環境變數",            ACCENT_BLUE),
        ("→ Factor VI：外部化 Session 狀態",         ACCENT_AMBER),
        ("→ Factor XI：Logs 輸出到 stdout",          ACCENT_GREEN),
        ("→ Factor V：Build-Release-Run 分離",       PART_COLORS[3]),
    ]

    y_start = Inches(1.6)
    card_h = Inches(0.75)
    y_spacing = Inches(0.82)

    for i, (text, accent) in enumerate(left_items):
        y = y_start + i * y_spacing
        add_card(slide, Inches(0.4), y, Inches(5.0), card_h,
                 body_lines=[text], accent=accent)

    for i, (text, accent) in enumerate(right_items):
        y = y_start + i * y_spacing
        add_card(slide, Inches(7.5), y, Inches(5.0), card_h,
                 body_lines=[text], accent=accent)

    # Center connecting column
    for i in range(6):
        y = y_start + i * y_spacing
        add_text(slide, "←→", Inches(5.6), y + Inches(0.2), Inches(1.6), Inches(0.4),
                 font_size=Pt(16), color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # Bottom summary card
    add_card(slide, Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.6),
             body_lines=["Container 天生符合大多數 12-Factor 原則 → Container + 12-Factor = 真正 Cloud-Ready 的應用"],
             accent=PART_COLORS[4])

    return slide


def slide_35(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "DevOps：打破 Dev 和 Ops 的高牆")
    add_part_label(slide, 5, "DevOps 大規模整合")
    add_page_number(slide, 35)

    # Left card - traditional silo culture
    add_card(slide, Inches(0.4), Inches(1.6), Inches(5.8), Inches(3.5),
             title="❌ 傳統孤島文化",
             body_lines=["Dev 團隊：「功能寫完了，丟給 Ops」",
                         "Ops 團隊：「怎麼又出問題了！」",
                         "",
                         "問題根源：",
                         "• Dev 不了解生產環境限制",
                         "• Ops 不了解程式碼變更影響",
                         "• 出問題互相推卸責任",
                         "• 部署頻率低，每次都是大事件"],
             accent=RGBColor(0xFF, 0x4A, 0x4A))

    # Right card - DevOps culture
    add_card(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(3.5),
             title="✅ DevOps 文化",
             body_lines=["共同擁有系統，共同負責品質",
                         "",
                         "核心原則：",
                         "• Dev 和 Ops 坐在同一團隊",
                         "• 開發者要對 On-call 負責",
                         "• 你 Build 它，你 Run 它",
                         "• 頻繁小批次部署，降低風險",
                         "• 失敗是學習機會，不是懲罰對象"],
             accent=ACCENT_GREEN)

    # Bottom summary card
    add_card(slide, Inches(0.4), Inches(5.35), Inches(12.5), Inches(1.0),
             title="DevOps 帶來的改變",
             body_lines=["部署頻率：數月一次 → 數十次/天  ·  MTTR：數天 → 數小時  ·  變更失敗率：大幅下降"],
             accent=PART_COLORS[5])

    add_callout(slide, "DevOps 是文化變革，工具是結果，不是原因。沒有文化基礎，買再多工具都沒用",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_36(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "CI/CD Pipeline：讓部署成為日常")
    add_part_label(slide, 5, "DevOps 大規模整合")
    add_page_number(slide, 36)

    # Flow diagram cards
    stages = [
        ("① Code Push", ACCENT_BLUE, ["git push", "PR Review", "Merge"]),
        ("② CI Pipeline", PART_COLORS[2], ["Build", "Lint", "Test"]),
        ("③ Image Build", PART_COLORS[3], ["docker build", "push to", "Registry"]),
        ("④ Staging Deploy", ACCENT_AMBER, ["自動部署", "驗收測試", "QA 確認"]),
        ("⑤ Prod Deploy", PART_COLORS[5], ["審核通過", "自動/手動", "Deploy"]),
    ]
    x_positions = [0.3, 3.0, 5.7, 8.4, 10.8]
    card_w = Inches(2.2)
    card_h = Inches(1.2)
    y = Inches(1.8)

    for i, (title, accent, body) in enumerate(stages):
        x = Inches(x_positions[i])
        add_card(slide, x, y, card_w, card_h, title=title, body_lines=body, accent=accent)
        if i < len(stages) - 1:
            arrow_x = x + card_w
            add_text(slide, "→", arrow_x, Inches(2.2), Inches(0.6), Inches(0.4),
                     font_size=Pt(18), color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # Two side-by-side cards below
    add_card(slide, Inches(0.4), Inches(3.2), Inches(5.8), Inches(2.5),
             title="沒有 CI/CD 的世界",
             body_lines=["部署是「重大事件」，半夜進行",
                         "手動執行 15 個步驟",
                         "出錯就怪「那個改程式的人」",
                         "一個月才部署一次，積累大量變更",
                         "每次部署都是在賭博"],
             accent=RGBColor(0xFF, 0x4A, 0x4A))

    add_card(slide, Inches(6.9), Inches(3.2), Inches(5.8), Inches(2.5),
             title="有 CI/CD 的世界",
             body_lines=["部署是「普通操作」，白天隨時進行",
                         "git push 後自動完成所有步驟",
                         "問題小且早發現，容易定位根因",
                         "每天可以部署多次，每次變更小",
                         "部署是工程師最無聊的工作"],
             accent=ACCENT_GREEN)

    add_callout(slide, "CI/CD 讓「部署」從「重大事件」變成「普通操作」，是 DevOps 文化的技術體現",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_37(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "大規模部署策略：如何安全上線？")
    add_part_label(slide, 5, "DevOps 大規模整合")
    add_page_number(slide, 37)

    strategies = [
        ("Rolling Update  滾動更新", ACCENT_BLUE,
         ["逐台更新，零停機",
          "",
          "流程：",
          "  Server 1 → 新版本 ✅",
          "  Server 2 → 新版本 ✅",
          "  Server 3 → 新版本 ✅",
          "",
          "✅ 零停機，資源不加倍",
          "✅ 自動，K8s 原生支援",
          "⚠ 瞬間新舊版本共存",
          "⚠ 回滾需要時間",
          "",
          "適用：大多數場景的首選"]),
        ("Blue-Green  藍綠部署", ACCENT_AMBER,
         ["兩套環境切換，零風險回滾",
          "",
          "Blue = 當前生產 (v1)",
          "Green = 新版本 (v2)",
          "",
          "1. Deploy v2 到 Green",
          "2. 驗證 Green 正常",
          "3. LB 切換到 Green",
          "4. Blue 保留備用",
          "",
          "✅ 回滾只需切換 LB (秒級)",
          "⚠ 需要 2x 資源成本"]),
        ("Canary Release  金絲雀", PART_COLORS[5],
         ["少量流量先驗證新版本",
          "",
          "流程：",
          "  5% 用戶 → v2 (新版本)",
          "  95% 用戶 → v1 (舊版本)",
          "  → 監控指標無異常",
          "  → 逐步 10% → 50% → 100%",
          "",
          "✅ 風險最低，影響面最小",
          "✅ 可以做 A/B Testing",
          "⚠ 需要流量分割能力",
          "",
          "適用：高風險功能上線"]),
    ]

    x_positions = [0.3, 4.6, 8.9]
    for i, (title, accent, body) in enumerate(strategies):
        add_card(slide, Inches(x_positions[i]), Inches(1.55), Inches(4.0), Inches(4.5),
                 title=title, body_lines=body, accent=accent)

    add_callout(slide, "選擇策略看業務容忍度：回滾速度優先 → Blue-Green；風險控制優先 → Canary",
                Inches(0.4), Inches(6.25), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_38(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "環境分層設計：Dev / Staging / Production")
    add_part_label(slide, 5, "DevOps 大規模整合")
    add_page_number(slide, 38)

    envs = [
        ("Development 環境", ACCENT_BLUE,
         ["用途：開發者本地開發",
          "使用者：工程師個人",
          "",
          "資料：假資料 / mock",
          "更新：隨時，無限制",
          "Config：.env.local",
          "",
          "工具：",
          "  Docker Compose 本地啟動",
          "  Hot reload 即時更新",
          "  Debug 模式開啟",
          "",
          "12-Factor：Factor X (Dev/Prod Parity)"]),
        ("Staging 環境", ACCENT_AMBER,
         ["用途：上線前最終驗證",
          "使用者：QA / PM / 利害關係人",
          "",
          "資料：生產資料的匿名副本",
          "更新：PR Merge 後自動部署",
          "Config：接近生產的設定",
          "",
          "原則：",
          "  Staging ≈ Production",
          "  越像，意外越少",
          "  完整流程測試（含第三方）",
          "",
          "12-Factor：Factor X (Dev/Prod Parity)"]),
        ("Production 環境", PART_COLORS[5],
         ["用途：真實用戶使用",
          "使用者：終端用戶",
          "",
          "資料：真實生產資料",
          "更新：通過 Staging 驗證後",
          "Config：生產 Secret，嚴格管控",
          "",
          "原則：",
          "  任何變更都有審核",
          "  有回滾計畫才能部署",
          "  監控全開，Alert 設定好",
          "",
          "Cost：最高，需要 HA 設計"]),
    ]

    x_positions = [0.3, 4.6, 8.9]
    for i, (title, accent, body) in enumerate(envs):
        add_card(slide, Inches(x_positions[i]), Inches(1.55), Inches(4.0), Inches(4.5),
                 title=title, body_lines=body, accent=accent)

    add_callout(slide, "Staging 與 Production 越相似，上線意外越少。最貴的是讓 Staging 和 Prod 不一樣的代價",
                Inches(0.4), Inches(6.25), Inches(12.5), Inches(0.55), style="warning")
    return slide


def slide_39(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "數十人團隊如何並行前進？")
    add_part_label(slide, 5, "DevOps 大規模整合")
    add_page_number(slide, 39)

    # Left column
    add_card(slide, Inches(0.4), Inches(1.6), Inches(5.5), Inches(2.5),
             title="❌ 常見問題",
             body_lines=["多團隊同時改同個 Repo → 互相衝突",
                         "部署要協調 → 誰先誰後？",
                         "一個 Bug 擋住所有人的上線",
                         "「等我這個功能合進去再說」"],
             accent=RGBColor(0xFF, 0x4A, 0x4A))

    add_card(slide, Inches(0.4), Inches(4.35), Inches(5.5), Inches(1.8),
             title="Monorepo vs Polyrepo",
             body_lines=["Monorepo：一個 Repo 放所有服務",
                         "  ✅ 跨服務修改方便  ⚠ 複雜度高",
                         "Polyrepo：每個服務獨立 Repo",
                         "  ✅ 團隊完全自治  ⚠ 跨服務協調難"],
             accent=ACCENT_AMBER)

    # Right column
    add_card(slide, Inches(6.2), Inches(1.6), Inches(6.5), Inches(3.2),
             title="✅ 服務邊界切分原則",
             body_lines=["一個團隊擁有一個（或少數幾個）服務",
                         "服務之間透過 API 溝通，不共享 DB",
                         "每個服務可以獨立部署，不依賴其他服務",
                         "",
                         "Conway's Law：",
                         "「系統架構會長得像組織架構」",
                         "→ 先設計好團隊結構，架構才會清晰"],
             accent=ACCENT_GREEN)

    add_card(slide, Inches(6.2), Inches(5.05), Inches(6.5), Inches(1.3),
             title="實務建議",
             body_lines=["Platform Team 負責共用基礎設施 (CI/CD / 監控)",
                         "Product Team 負責各自的服務，自主部署"],
             accent=PART_COLORS[5])

    add_callout(slide, "讓每個團隊能獨立 Deploy，是大規模組織前進速度的關鍵",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_40(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "GitOps：以 Git 為唯一事實來源")
    add_part_label(slide, 5, "DevOps 大規模整合")
    add_page_number(slide, 40)

    # Top flow
    flow_stages = [
        ("Developer", ACCENT_BLUE, ["寫 Code", "開 PR"]),
        ("PR Review", PART_COLORS[2], ["Code Review", "Approved"]),
        ("Git Merge", ACCENT_AMBER, ["Merge to", "main branch"]),
        ("CD Tool", PART_COLORS[3], ["ArgoCD", "Flux CD"]),
        ("Cluster", PART_COLORS[5], ["K8s / Prod", "自動同步"]),
    ]
    x_positions = [0.3, 3.0, 5.7, 8.4, 10.8]
    card_w = Inches(2.2)
    card_h = Inches(1.0)
    y = Inches(1.6)

    for i, (title, accent, body) in enumerate(flow_stages):
        x = Inches(x_positions[i])
        add_card(slide, x, y, card_w, card_h, title=title, body_lines=body, accent=accent)
        if i < len(flow_stages) - 1:
            arrow_x = x + card_w
            add_text(slide, "→", arrow_x, Inches(2.0), Inches(0.6), Inches(0.4),
                     font_size=Pt(18), color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # Middle principle cards
    principles = [
        ("① 宣告式設定", ACCENT_BLUE,
         ["所有部署設定都在 Git Repo", "YAML 描述「期望狀態」", "不是「怎麼做」，而是「要什麼」"]),
        ("② PR-Driven 部署", PART_COLORS[2],
         ["任何部署變更都要開 PR", "有 Review 才能上線", "Git history = 完整部署紀錄"]),
        ("③ 自動同步", PART_COLORS[5],
         ["CD 工具持續監控 Git", "偵測到變更自動部署", "系統自動對齊 Git 狀態"]),
    ]
    x_positions_mid = [0.3, 4.6, 8.9]
    for i, (title, accent, body) in enumerate(principles):
        add_card(slide, Inches(x_positions_mid[i]), Inches(2.9), Inches(4.0), Inches(2.0),
                 title=title, body_lines=body, accent=accent)

    # Bottom summary card
    add_card(slide, Inches(0.4), Inches(5.1), Inches(12.5), Inches(1.1),
             title="GitOps 的好處",
             body_lines=["可審計：每次部署都有 Git commit  ·  可追蹤：誰在什麼時間改了什麼  ·  可回滾：git revert 即可回到任何版本"],
             accent=ACCENT_GREEN)

    add_callout(slide, "Git commit history = 完整部署紀錄。任何人任何時候都知道生產環境是什麼狀態",
                Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_41(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "Feature Flags：解耦 Deploy 與 Release")
    add_part_label(slide, 5, "DevOps 大規模整合")
    add_page_number(slide, 41)

    # Left card
    add_card(slide, Inches(0.4), Inches(1.6), Inches(5.8), Inches(2.8),
             title="❌ 沒有 Feature Flags",
             body_lines=["功能寫完才能部署",
                         "→ 分支長期存在 → Merge 衝突",
                         "→ 部署和上線綁死在一起",
                         "→ 新功能有 Bug？整個版本回滾",
                         "→ A/B Testing？重新部署兩個版本"],
             accent=RGBColor(0xFF, 0x4A, 0x4A))

    # Right card
    add_card(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(2.8),
             title="✅ 有 Feature Flags",
             body_lines=["if feature_flags.is_enabled('new_checkout'):",
                         "    show_new_checkout_ui()",
                         "else:",
                         "    show_old_checkout_ui()",
                         "",
                         "→ 隨時可以部署，用 Flag 控制開放"],
             accent=ACCENT_GREEN)

    # Use case cards
    use_cases = [
        ("Gradual Rollout", ACCENT_BLUE, ["先開放 1% → 5% → 100%", "確認無問題再全量"]),
        ("A/B Testing", PART_COLORS[2], ["A 組用舊版，B 組用新版", "用數據決定哪個更好"]),
        ("Kill Switch", ACCENT_AMBER, ["出問題即時關閉功能", "不需要 Rollback 整個部署"]),
    ]
    x_positions = [0.3, 4.6, 8.9]
    for i, (title, accent, body) in enumerate(use_cases):
        add_card(slide, Inches(x_positions[i]), Inches(4.65), Inches(3.9), Inches(1.8),
                 title=title, body_lines=body, accent=accent)

    add_callout(slide, "Deploy ≠ Release。部署是技術動作，上線是商業決策。Feature Flags 讓兩者解耦",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55), style="tip")
    return slide


def slide_42(prs):
    slide = new_slide(prs)
    add_slide_title(slide, "Part 5 小結：DevOps 讓數十人團隊安全前進")
    add_part_label(slide, 5, "DevOps 大規模整合")
    add_page_number(slide, 42)

    # Integrated flow diagram
    flow_items = [
        ("DevOps 文化", PART_COLORS[5]),
        ("CI/CD", PART_COLORS[2]),
        ("環境分層", ACCENT_AMBER),
        ("部署策略", ACCENT_BLUE),
        ("GitOps", PART_COLORS[3]),
        ("Feature Flags", ACCENT_GREEN),
        ("可靠上線", PART_COLORS[5]),
    ]
    x_positions = [0.3, 2.2, 4.1, 6.0, 7.9, 9.8, 11.5]
    card_w = Inches(1.6)
    card_h = Inches(1.0)
    y = Inches(1.7)

    for i, (label, accent) in enumerate(flow_items):
        x = Inches(x_positions[i])
        add_card(slide, x, y, card_w, card_h, body_lines=[label], accent=accent)
        if i < len(flow_items) - 1:
            arrow_x = x + card_w
            add_text(slide, "→", arrow_x, Inches(2.0), Inches(0.4), Inches(0.35),
                     font_size=Pt(16), color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # Key insight card
    add_card(slide, Inches(0.4), Inches(3.0), Inches(12.5), Inches(1.5),
             title="核心認知",
             body_lines=["CI/CD 不是工具，是文化和流程",
                         "讓每個 PR Merge 都可以安全上線，讓每個 Deploy 都可以快速回滾",
                         "目標：讓部署成為「最無聊」的工作，工程師應該把精力放在「產品創新」"],
             accent=PART_COLORS[5])

    # Metric cards
    metrics = [
        ("部署頻率", ACCENT_BLUE, ["目標：每天多次", "指標：Deploy Frequency"]),
        ("變更前置時間", ACCENT_GREEN, ["Commit → Production", "指標：Lead Time for Change"]),
        ("平均恢復時間", ACCENT_AMBER, ["出事到恢復多快", "指標：MTTR"]),
    ]
    x_positions_m = [0.3, 4.6, 8.9]
    for i, (title, accent, body) in enumerate(metrics):
        add_card(slide, Inches(x_positions_m[i]), Inches(4.75), Inches(3.9), Inches(1.6),
                 title=title, body_lines=body, accent=accent)

    add_callout(slide, "下一步 → Part 6：整個 SDLC 閉環，從 Idea 到監控生產環境的完整旅程",
                Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.55), style="tip")
    return slide


if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_cover(prs)
    slide_02_agenda(prs)
    for fn in [slide_03, slide_04, slide_05, slide_06, slide_07,
               slide_08, slide_09, slide_10, slide_11, slide_12,
               slide_13, slide_14, slide_15, slide_16, slide_17,
               slide_18, slide_19, slide_20,
               slide_21, slide_22, slide_23, slide_24, slide_25, slide_26,
               slide_27, slide_28, slide_29, slide_30, slide_31,
               slide_32, slide_33, slide_34,
               slide_35, slide_36, slide_37, slide_38, slide_39, slide_40,
               slide_41, slide_42]:
        fn(prs)

    prs.save("cloud_native_slides_v2.pptx")
    print(f"✅ Generated {len(prs.slides)} slides → cloud_native_slides_v2.pptx")
