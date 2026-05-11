from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Theme ─────────────────────────────────────────────────────────────────────
BG_COLOR     = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT_BLUE  = RGBColor(0x00, 0xBF, 0xFF)
ACCENT_GREEN = RGBColor(0x39, 0xFF, 0x14)
ACCENT_AMBER = RGBColor(0xFF, 0xBF, 0x00)
TEXT_PRIMARY = RGBColor(0xE8, 0xED, 0xF2)
TEXT_DIM     = RGBColor(0x8F, 0xA3, 0xBF)
FONT_TITLE   = "JetBrains Mono"
FONT_BODY    = "Inter"
SLIDE_W      = Inches(13.33)
SLIDE_H      = Inches(7.5)

# DC accent colors
DC_BORDERS      = [ACCENT_BLUE, RGBColor(0x7B, 0x2F, 0xFF), ACCENT_GREEN]
DC_FILLS        = [RGBColor(0x06, 0x14, 0x26), RGBColor(0x10, 0x08, 0x22), RGBColor(0x06, 0x14, 0x10)]
DC_HEADER_FILLS = [RGBColor(0x00, 0x30, 0x70), RGBColor(0x30, 0x10, 0x70), RGBColor(0x00, 0x40, 0x20)]
DC_LABELS       = ["Data Center 1", "Data Center 2", "Data Center 3"]

RACK_FILL   = RGBColor(0x14, 0x20, 0x2E)
RACK_BORDER = RGBColor(0x20, 0x35, 0x50)
NODE_FILL   = RGBColor(0x00, 0x8A, 0xC0)
NODE_BORDER = RGBColor(0x00, 0xC0, 0xFF)

# Applications shown running on K8s
APPS = [
    ("Frontend",     RGBColor(0x00, 0x50, 0x90)),
    ("Backend API",  RGBColor(0x00, 0x50, 0x90)),
    ("Worker",       RGBColor(0x00, 0x50, 0x90)),
    ("Redis Cache",  RGBColor(0x50, 0x20, 0x00)),
    ("PostgreSQL",   RGBColor(0x50, 0x20, 0x00)),
    ("Monitoring",   RGBColor(0x10, 0x40, 0x10)),
]


def create_k8s_intro():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR

    # ── Helpers ───────────────────────────────────────────────────────
    def R(l, t, w, h, fill, border=None, bw=Pt(1)):
        s = slide.shapes.add_shape(1, l, t, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        if border:
            s.line.color.rgb = border
            s.line.width = bw
        else:
            s.line.fill.background()
        return s

    def T(txt, l, t, w, h, fn=FONT_BODY, fs=Pt(12),
          color=TEXT_PRIMARY, bold=False, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = txt
        run.font.name = fn
        run.font.size = fs
        run.font.color.rgb = color
        run.font.bold = bold
        return tb

    # ── Top accent bar ────────────────────────────────────────────────
    R(Inches(0), Inches(0), SLIDE_W, Inches(0.05), ACCENT_BLUE)

    # ── Title ─────────────────────────────────────────────────────────
    T("Kubernetes 跨實體分散架構",
      Inches(0.4), Inches(0.08), Inches(10), Inches(0.52),
      fn=FONT_TITLE, fs=Pt(28), color=ACCENT_BLUE, bold=True)
    T("分散到不同實體位置的大運算平台，任何單一實體單元損毀不會讓叢集整體損毀",
      Inches(0.4), Inches(0.62), Inches(12.5), Inches(0.28),
      fs=Pt(12), color=TEXT_DIM)

    # ── Applications Layer ────────────────────────────────────────────
    APP_Y  = Inches(0.96)
    APP_H  = Inches(1.25)
    APP_W  = Inches(1.88)
    APP_GAP = Inches(0.17)

    # Layer label
    T("Applications  (運行於 Kubernetes 之上)",
      Inches(0.3), APP_Y - Inches(0.01), Inches(5), Inches(0.25),
      fn=FONT_TITLE, fs=Pt(9), color=TEXT_DIM)

    for i, (name, fill) in enumerate(APPS):
        ax = Inches(0.3) + i * (APP_W + APP_GAP)
        # Card background
        R(ax, APP_Y + Inches(0.22), APP_W, APP_H, fill, ACCENT_BLUE, Pt(1.2))
        # Pod icon bar
        R(ax, APP_Y + Inches(0.22), APP_W, Inches(0.32), ACCENT_BLUE)
        T("⬡ Pod", ax + Inches(0.08), APP_Y + Inches(0.26),
          APP_W - Inches(0.16), Inches(0.24),
          fn=FONT_TITLE, fs=Pt(8), color=RGBColor(0x0D, 0x1B, 0x2A), bold=True)
        T(name, ax + Inches(0.1), APP_Y + Inches(0.62),
          APP_W - Inches(0.2), Inches(0.45),
          fn=FONT_TITLE, fs=Pt(11), color=TEXT_PRIMARY, bold=True,
          align=PP_ALIGN.CENTER)
        T("Deployment", ax + Inches(0.1), APP_Y + Inches(1.08),
          APP_W - Inches(0.2), Inches(0.25),
          fs=Pt(9), color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # Connector arrows Apps → K8s
    T("↕  K8s 調度與管理",
      Inches(5.5), APP_Y + APP_H + Inches(0.22), Inches(3), Inches(0.28),
      fn=FONT_TITLE, fs=Pt(9), color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # ── Kubernetes Cluster bar ────────────────────────────────────────
    CLUSTER_Y = APP_Y + APP_H + Inches(0.56)
    CLUSTER_H = Inches(0.42)
    R(Inches(0.3), CLUSTER_Y, Inches(12.73), CLUSTER_H,
      RGBColor(0x00, 0x28, 0x55), ACCENT_BLUE, Pt(2))
    T("⎈  Kubernetes Cluster",
      Inches(0.45), CLUSTER_Y + Inches(0.06), Inches(5), Inches(0.34),
      fn=FONT_TITLE, fs=Pt(13), color=ACCENT_BLUE, bold=True)
    T("Control Plane:  API Server  /  etcd  /  Scheduler  /  Controller Manager",
      Inches(5.5), CLUSTER_Y + Inches(0.08), Inches(7.2), Inches(0.28),
      fs=Pt(10), color=TEXT_DIM, align=PP_ALIGN.RIGHT)

    # Connector arrows K8s → Physical
    T("↕  實際運行在底層實體機器上",
      Inches(4.5), CLUSTER_Y + CLUSTER_H + Inches(0.06), Inches(4), Inches(0.28),
      fn=FONT_TITLE, fs=Pt(9), color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # ── Physical Infrastructure (simplified — no Room layer) ──────────
    DC_Y  = CLUSTER_Y + CLUSTER_H + Inches(0.38)
    DC_H  = Inches(2.72)
    DC_W  = Inches(4.11)
    DC_GAP = Inches(0.2)
    DC_XS = [
        Inches(0.3),
        Inches(0.3) + DC_W + DC_GAP,
        Inches(0.3) + 2 * (DC_W + DC_GAP),
    ]

    RACK_PAD    = Inches(0.09)
    RACK_W      = (DC_W - 2 * RACK_PAD - 3 * RACK_PAD) / 2   # 2 racks + padding
    RACK_H      = DC_H - Inches(0.55)
    RACK_LABELS = ["Rack 1", "Rack 2"]
    NODE_W      = RACK_W - Inches(0.12)
    NODE_H      = Inches(0.36)

    for di in range(3):
        dx = DC_XS[di]

        # DC outer box
        R(dx, DC_Y, DC_W, DC_H, DC_FILLS[di], DC_BORDERS[di], Pt(2))

        # DC header strip
        R(dx, DC_Y, DC_W, Inches(0.38), DC_HEADER_FILLS[di])
        T(DC_LABELS[di],
          dx + Inches(0.12), DC_Y + Inches(0.05), DC_W - Inches(0.24), Inches(0.30),
          fn=FONT_TITLE, fs=Pt(11), color=TEXT_PRIMARY, bold=True)

        for rki in range(2):
            rkx = dx + RACK_PAD + rki * (RACK_W + RACK_PAD * 2)
            rky = DC_Y + Inches(0.46)

            # Rack box
            R(rkx, rky, RACK_W, RACK_H, RACK_FILL, RACK_BORDER, Pt(1))
            T(RACK_LABELS[rki],
              rkx + Inches(0.06), rky + Inches(0.05),
              RACK_W - Inches(0.12), Inches(0.24),
              fn=FONT_TITLE, fs=Pt(9), color=TEXT_DIM)

            # 3 server nodes per rack
            node_x = rkx + Inches(0.06)
            for ni in range(3):
                ny = rky + Inches(0.33) + ni * Inches(0.44)
                node_num = di * 6 + rki * 3 + ni + 1
                R(node_x, ny, NODE_W, NODE_H, NODE_FILL, NODE_BORDER, Pt(0.8))
                T(f"Node {node_num:02d}",
                  node_x + Inches(0.06), ny + Inches(0.05),
                  NODE_W - Inches(0.12), NODE_H - Inches(0.08),
                  fs=Pt(8), color=RGBColor(0x0D, 0x1B, 0x2A), bold=True)

    # ── Callout ───────────────────────────────────────────────────────
    CALLOUT_Y = DC_Y + DC_H + Inches(0.10)
    CALLOUT_H = Inches(0.52)
    R(Inches(0.3), CALLOUT_Y, Inches(12.73), CALLOUT_H,
      RGBColor(0x00, 0x18, 0x08), ACCENT_GREEN, Pt(1.5))
    T("💡  任何單一 DC / Rack 損毀，Kubernetes 仍可透過其他節點繼續調度應用"
      " — 上層 Application 對底層實體故障無感",
      Inches(0.45), CALLOUT_Y + Inches(0.09),
      Inches(12.3), CALLOUT_H - Inches(0.14),
      fs=Pt(11.5), color=TEXT_PRIMARY)

    # ── Bottom accent bar ─────────────────────────────────────────────
    R(Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), ACCENT_BLUE)

    prs.save("k8s_introduction.pptx")
    print("✅ Generated k8s_introduction.pptx (1 slide)")


create_k8s_intro()
